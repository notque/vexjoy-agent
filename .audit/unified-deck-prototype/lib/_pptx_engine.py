#!/usr/bin/env python3
"""
Deterministic PPTX generator.

Reads a slide map JSON and design config JSON, produces a .pptx file.
No LLM calls -- pure mechanical slide construction using python-pptx.

Usage:
    python3 generate_pptx.py --slide-map slides.json --design design.json --output deck.pptx

Exit codes:
    0 = success
    1 = missing dependencies
    2 = invalid input
    3 = generation failed
"""

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


# ---------------------------------------------------------------------------
# Palette definitions
# ---------------------------------------------------------------------------
PALETTES = {
    "corporate": {
        "primary": "#2C3E50",
        "secondary": "#34495E",
        "accent": "#E74C3C",
        "background": "#FFFFFF",
        "text": "#2C3E50",
        "muted": "#BDC3C7",
    },
    "tech": {
        "primary": "#1A1A2E",
        "secondary": "#16213E",
        "accent": "#0F3460",
        "background": "#F5F5F5",
        "text": "#1A1A2E",
        "muted": "#A0A0A0",
    },
    "warm": {
        "primary": "#D4A574",
        "secondary": "#C68B59",
        "accent": "#8B4513",
        "background": "#FFF8F0",
        "text": "#3D2B1F",
        "muted": "#D2B48C",
    },
    "ocean": {
        "primary": "#006994",
        "secondary": "#008B8B",
        "accent": "#20B2AA",
        "background": "#F0FFFF",
        "text": "#003333",
        "muted": "#B0C4DE",
    },
    "midnight": {
        "primary": "#1B1464",
        "secondary": "#2E1A47",
        "accent": "#7B2FBE",
        "background": "#0D0D0D",
        "text": "#E0E0E0",
        "muted": "#4A4A6A",
    },
    "forest": {
        "primary": "#2D5016",
        "secondary": "#3A6B1E",
        "accent": "#7CB342",
        "background": "#F5F9F0",
        "text": "#1A2E0A",
        "muted": "#A5D6A7",
    },
    "sunset": {
        "primary": "#FF6B35",
        "secondary": "#F7931E",
        "accent": "#FFD700",
        "background": "#FFF5E6",
        "text": "#333333",
        "muted": "#FFE0B2",
    },
    "minimal": {
        "primary": "#333333",
        "secondary": "#666666",
        "accent": "#0066CC",
        "background": "#FFFFFF",
        "text": "#333333",
        "muted": "#CCCCCC",
    },
}


def hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert '#RRGGBB' to RGBColor."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def get_palette(name: str) -> dict:
    """Return palette dict with RGBColor values."""
    raw = PALETTES.get(name.lower(), PALETTES["minimal"])
    return {role: hex_to_rgb(color) for role, color in raw.items()}


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def set_slide_background(slide, color: RGBColor):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    bold=False,
    italic=False,
    color=None,
    alignment=None,
    word_wrap=True,
    line_spacing=None,
):
    """Add a text box with formatted text to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    if color:
        p.font.color.rgb = color
    if alignment:
        p.alignment = alignment
    if line_spacing:
        p.line_spacing = line_spacing
    p.font.name = "Calibri"
    return txBox


def add_bullet_list(slide, left, top, width, height, bullets, font_size=18, color=None, line_spacing=1.15):
    """Add a bulleted list text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if isinstance(bullet, dict):
            text = bullet.get("text", "")
            level = bullet.get("level", 0)
        else:
            text = str(bullet)
            level = 0

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = text
        p.font.size = Pt(font_size)
        p.font.name = "Calibri"
        p.level = level
        p.space_after = Pt(6)
        p.line_spacing = line_spacing
        if color:
            p.font.color.rgb = color

    return txBox


def build_title_slide(prs, slide_data, palette):
    """Build a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide, palette["background"])

    title = slide_data.get("title", "Untitled Presentation")
    subtitle = slide_data.get("subtitle", "")

    add_text_box(
        slide,
        1,
        2,
        11.333,
        2,
        title,
        font_size=44,
        bold=True,
        color=palette["text"],
        alignment=PP_ALIGN.CENTER,
    )

    if subtitle:
        add_text_box(
            slide,
            2,
            4.2,
            9.333,
            1,
            subtitle,
            font_size=22,
            color=palette["secondary"],
            alignment=PP_ALIGN.CENTER,
        )

    return slide


def build_section_divider(prs, slide_data, palette):
    """Build a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, palette["primary"])

    title = slide_data.get("title", "Section")
    subtitle = slide_data.get("subtitle", "")

    # White text on primary background
    white = RGBColor(0xFF, 0xFF, 0xFF)
    add_text_box(
        slide,
        1,
        2.5,
        11.333,
        2,
        title,
        font_size=36,
        bold=True,
        color=white,
        alignment=PP_ALIGN.LEFT,
    )

    if subtitle:
        add_text_box(
            slide,
            1,
            4.5,
            11.333,
            1,
            subtitle,
            font_size=20,
            color=RGBColor(0xE0, 0xE0, 0xE0),
            alignment=PP_ALIGN.LEFT,
        )

    return slide


def build_content_bullets(prs, slide_data, palette):
    """Build a standard content slide with headline and bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, palette["background"])

    title = slide_data.get("title", "")
    bullets = slide_data.get("bullets", [])

    if title:
        add_text_box(
            slide,
            0.5,
            0.5,
            12.333,
            1,
            title,
            font_size=28,
            bold=True,
            color=palette["text"],
        )

    if bullets:
        add_bullet_list(
            slide,
            0.5,
            1.8,
            12.333,
            5.2,
            bullets,
            font_size=18,
            color=palette["text"],
        )

    return slide


def build_two_column(prs, slide_data, palette):
    """Build a two-column content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, palette["background"])

    title = slide_data.get("title", "")
    left_content = slide_data.get("left", {})
    right_content = slide_data.get("right", {})

    if title:
        add_text_box(
            slide,
            0.5,
            0.5,
            12.333,
            1,
            title,
            font_size=28,
            bold=True,
            color=palette["text"],
        )

    # Left column header
    left_header = left_content.get("header", "")
    if left_header:
        add_text_box(
            slide,
            0.5,
            1.8,
            5.9,
            0.6,
            left_header,
            font_size=20,
            bold=True,
            color=palette["primary"],
        )

    left_bullets = left_content.get("bullets", [])
    if left_bullets:
        top = 2.5 if left_header else 1.8
        add_bullet_list(
            slide,
            0.5,
            top,
            5.9,
            4.5,
            left_bullets,
            font_size=18,
            color=palette["text"],
        )

    # Right column header
    right_header = right_content.get("header", "")
    if right_header:
        add_text_box(
            slide,
            6.933,
            1.8,
            5.9,
            0.6,
            right_header,
            font_size=20,
            bold=True,
            color=palette["primary"],
        )

    right_bullets = right_content.get("bullets", [])
    if right_bullets:
        top = 2.5 if right_header else 1.8
        add_bullet_list(
            slide,
            6.933,
            top,
            5.9,
            4.5,
            right_bullets,
            font_size=18,
            color=palette["text"],
        )

    return slide


def build_quote(prs, slide_data, palette):
    """Build a quote/callout slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, palette["background"])

    quote_text = slide_data.get("quote", "")
    attribution = slide_data.get("attribution", "")

    if quote_text:
        add_text_box(
            slide,
            1.5,
            2,
            10.333,
            3,
            f'"{quote_text}"',
            font_size=28,
            italic=True,
            color=palette["text"],
            alignment=PP_ALIGN.CENTER,
        )

    if attribution:
        add_text_box(
            slide,
            2,
            5.2,
            9.333,
            0.8,
            f"-- {attribution}",
            font_size=16,
            color=palette["muted"],
            alignment=PP_ALIGN.CENTER,
        )

    return slide


def build_table_slide(prs, slide_data, palette):
    """Build a table slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, palette["background"])

    title = slide_data.get("title", "")
    headers = slide_data.get("headers", [])
    rows_data = slide_data.get("rows", [])

    if title:
        add_text_box(
            slide,
            0.5,
            0.5,
            12.333,
            1,
            title,
            font_size=28,
            bold=True,
            color=palette["text"],
        )

    if headers and rows_data:
        num_rows = len(rows_data) + 1  # +1 for header
        num_cols = len(headers)

        table_shape = slide.shapes.add_table(
            num_rows,
            num_cols,
            Inches(0.5),
            Inches(1.8),
            Inches(12.333),
            Inches(min(5.0, 0.5 * num_rows + 0.5)),
        )
        table = table_shape.table

        # Header row
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = palette["primary"]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                paragraph.font.bold = True
                paragraph.font.size = Pt(14)
                paragraph.font.name = "Calibri"

        # Data rows
        for row_idx, row in enumerate(rows_data):
            for col_idx, cell_text in enumerate(row):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_text)
                # Alternating row colors
                if row_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = palette.get("muted", RGBColor(0xF0, 0xF0, 0xF0))
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.font.name = "Calibri"
                    paragraph.font.color.rgb = palette["text"]

    return slide


def build_image_text(prs, slide_data, palette):
    """Build an image + text slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, palette["background"])

    title = slide_data.get("title", "")
    image_path = slide_data.get("image_path", "")
    bullets = slide_data.get("bullets", [])
    image_side = slide_data.get("image_side", "left")

    if title:
        add_text_box(
            slide,
            0.5,
            0.5,
            12.333,
            1,
            title,
            font_size=28,
            bold=True,
            color=palette["text"],
        )

    if image_side == "left":
        img_left, text_left = 0.5, 6.933
    else:
        img_left, text_left = 6.933, 0.5

    # Add image if it exists
    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(
            image_path,
            Inches(img_left),
            Inches(1.8),
            Inches(5.9),
            Inches(5.2),
        )

    # Add text bullets
    if bullets:
        add_bullet_list(
            slide,
            text_left,
            1.8,
            5.9,
            5.2,
            bullets,
            font_size=18,
            color=palette["text"],
        )

    return slide


def build_closing_slide(prs, slide_data, palette):
    """Build a closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, palette["primary"])

    main_text = slide_data.get("title", "Thank You")
    subtitle = slide_data.get("subtitle", "")

    white = RGBColor(0xFF, 0xFF, 0xFF)
    add_text_box(
        slide,
        1,
        2.5,
        11.333,
        2,
        main_text,
        font_size=36,
        bold=True,
        color=white,
        alignment=PP_ALIGN.CENTER,
    )

    if subtitle:
        add_text_box(
            slide,
            2,
            4.8,
            9.333,
            1,
            subtitle,
            font_size=18,
            color=RGBColor(0xE0, 0xE0, 0xE0),
            alignment=PP_ALIGN.CENTER,
        )

    return slide


# ---------------------------------------------------------------------------
# Layout dispatcher
# ---------------------------------------------------------------------------

LAYOUT_BUILDERS = {
    "title": build_title_slide,
    "section": build_section_divider,
    "section_divider": build_section_divider,
    "content": build_content_bullets,
    "bullets": build_content_bullets,
    "content_bullets": build_content_bullets,
    "two_column": build_two_column,
    "two-column": build_two_column,
    "quote": build_quote,
    "callout": build_quote,
    "table": build_table_slide,
    "image_text": build_image_text,
    "image-text": build_image_text,
    "closing": build_closing_slide,
}


def build_presentation(slide_map: list, design: dict, output_path: str) -> str:
    """Build a complete presentation from slide map and design config.

    Args:
        slide_map: List of slide dicts with 'type' and content fields.
        design: Dict with 'palette' name and optional overrides.
        output_path: Path to write the .pptx file.

    Returns:
        Path to the created .pptx file.
    """
    palette_name = design.get("palette", "minimal")
    palette = get_palette(palette_name)

    # Support template-based generation
    template_path = design.get("template_path")
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    # Set widescreen (16:9) dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in slide_map:
        slide_type = slide_data.get("type", "content").lower().replace(" ", "_")
        builder = LAYOUT_BUILDERS.get(slide_type, build_content_bullets)
        builder(prs, slide_data, palette)

    # Ensure output directory exists
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    prs.save(str(output))
    return str(output)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------


def main():
    parser = argparse.ArgumentParser(description="Generate a PPTX presentation from a slide map JSON.")
    parser.add_argument(
        "--slide-map",
        required=True,
        help="Path to slide map JSON file",
    )
    parser.add_argument(
        "--design",
        required=True,
        help="Path to design config JSON file",
    )
    parser.add_argument(
        "--output",
        required=True,
        help="Output .pptx file path",
    )
    args = parser.parse_args()

    # Validate inputs
    slide_map_path = Path(args.slide_map)
    design_path = Path(args.design)

    if not slide_map_path.exists():
        print(f"ERROR: Slide map not found: {slide_map_path}", file=sys.stderr)
        sys.exit(2)

    if not design_path.exists():
        print(f"ERROR: Design config not found: {design_path}", file=sys.stderr)
        sys.exit(2)

    try:
        with open(slide_map_path) as f:
            slide_map = json.load(f)
    except json.JSONDecodeError as e:
        print(f"ERROR: Invalid slide map JSON: {e}", file=sys.stderr)
        sys.exit(2)

    try:
        with open(design_path) as f:
            design = json.load(f)
    except json.JSONDecodeError as e:
        print(f"ERROR: Invalid design config JSON: {e}", file=sys.stderr)
        sys.exit(2)

    if not isinstance(slide_map, list):
        print("ERROR: Slide map must be a JSON array of slide objects", file=sys.stderr)
        sys.exit(2)

    if len(slide_map) == 0:
        print("ERROR: Slide map is empty", file=sys.stderr)
        sys.exit(2)

    try:
        result = build_presentation(slide_map, design, args.output)
        file_size = Path(result).stat().st_size
        print(f"SUCCESS: Generated {len(slide_map)} slides")
        print(f"  Output: {result}")
        print(f"  Size: {file_size:,} bytes")
        print(f"  Palette: {design.get('palette', 'minimal')}")
    except Exception as e:
        print(f"ERROR: Generation failed: {e}", file=sys.stderr)
        sys.exit(3)


if __name__ == "__main__":
    main()
