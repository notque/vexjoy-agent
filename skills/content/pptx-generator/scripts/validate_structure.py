#!/usr/bin/env python3
"""
Validate PPTX structure: check slide count, titles, and basic integrity.

Usage:
    python3 validate_structure.py --input deck.pptx --expected-slides 10
    python3 validate_structure.py --input deck.pptx --slide-map slides.json

Exit codes:
    0 = validation passed
    1 = missing dependencies
    2 = validation failed
    3 = invalid input
"""

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


def extract_slide_info(pptx_path: str) -> list[dict]:
    """Extract basic info from each slide in a PPTX file.

    Returns list of dicts with 'index', 'title' (first text found),
    'shape_count', and 'text_content'.
    """
    prs = Presentation(pptx_path)
    slides_info = []

    for i, slide in enumerate(prs.slides):
        info = {
            "index": i + 1,
            "title": "",
            "shape_count": len(slide.shapes),
            "text_content": [],
        }

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    info["text_content"].append(text)
                    if not info["title"]:
                        info["title"] = text

        slides_info.append(info)

    return slides_info


def validate_slide_count(slides_info: list, expected: int) -> list[str]:
    """Check that slide count matches expected."""
    issues = []
    actual = len(slides_info)
    if actual != expected:
        issues.append(f"Slide count mismatch: expected {expected}, got {actual}")
    return issues


def validate_against_slide_map(slides_info: list, slide_map: list) -> list[str]:
    """Validate slides against the slide map.

    Checks:
    - Slide count matches
    - Each slide has some text content
    - Title slide exists (first slide)
    """
    issues = []

    # Check count
    if len(slides_info) != len(slide_map):
        issues.append(
            f"Slide count mismatch: slide map has {len(slide_map)} slides, PPTX has {len(slides_info)} slides"
        )

    # Check each slide has content
    for i, (actual, expected) in enumerate(zip(slides_info, slide_map, strict=False)):
        slide_num = i + 1
        expected_title = expected.get("title", "")

        if not actual["text_content"]:
            issues.append(f"Slide {slide_num}: no text content found")

        if expected_title and not any(expected_title.lower() in text.lower() for text in actual["text_content"]):
            issues.append(f"Slide {slide_num}: expected title '{expected_title}' not found in slide text")

    return issues


def validate_no_empty_slides(slides_info: list) -> list[str]:
    """Check for slides with zero shapes or no text."""
    issues = []
    for slide in slides_info:
        if slide["shape_count"] == 0:
            issues.append(f"Slide {slide['index']}: has zero shapes (empty slide)")
        elif not slide["text_content"]:
            # Slides with shapes but no text might be image-only, which is OK
            pass
    return issues


def validate_file_integrity(pptx_path: str) -> list[str]:
    """Basic file integrity checks."""
    issues = []
    path = Path(pptx_path)

    if not path.exists():
        issues.append(f"File does not exist: {pptx_path}")
        return issues

    if path.stat().st_size == 0:
        issues.append(f"File is empty (0 bytes): {pptx_path}")
        return issues

    if path.stat().st_size < 1000:
        issues.append(f"File suspiciously small ({path.stat().st_size} bytes): {pptx_path}")

    return issues


def validate(pptx_path: str, expected_slides: int | None = None, slide_map: list | None = None) -> dict:
    """Run all validations and return results.

    Returns dict with 'passed', 'issues', and 'slides_info'.
    """
    all_issues = []

    # File integrity
    file_issues = validate_file_integrity(pptx_path)
    all_issues.extend(file_issues)

    if file_issues:
        # Can't proceed with structural checks if file is bad
        return {
            "passed": False,
            "issues": all_issues,
            "slides_info": [],
        }

    # Extract slide info
    try:
        slides_info = extract_slide_info(pptx_path)
    except Exception as e:
        all_issues.append(f"Failed to read PPTX: {e}")
        return {
            "passed": False,
            "issues": all_issues,
            "slides_info": [],
        }

    # Slide count
    if expected_slides is not None:
        all_issues.extend(validate_slide_count(slides_info, expected_slides))

    # Slide map comparison
    if slide_map is not None:
        all_issues.extend(validate_against_slide_map(slides_info, slide_map))

    # Empty slides
    all_issues.extend(validate_no_empty_slides(slides_info))

    return {
        "passed": len(all_issues) == 0,
        "issues": all_issues,
        "slides_info": slides_info,
    }


def main():
    parser = argparse.ArgumentParser(description="Validate PPTX structure and content.")
    parser.add_argument(
        "--input",
        required=True,
        help="Path to .pptx file to validate",
    )
    parser.add_argument(
        "--expected-slides",
        type=int,
        default=None,
        help="Expected number of slides",
    )
    parser.add_argument(
        "--slide-map",
        default=None,
        help="Path to slide map JSON for content validation",
    )
    parser.add_argument(
        "--json",
        action="store_true",
        help="Output results as JSON",
    )
    args = parser.parse_args()

    # Load slide map if provided
    slide_map = None
    if args.slide_map:
        try:
            with open(args.slide_map) as f:
                slide_map = json.load(f)
        except (json.JSONDecodeError, FileNotFoundError) as e:
            print(f"ERROR: Invalid slide map: {e}", file=sys.stderr)
            sys.exit(3)

    result = validate(args.input, args.expected_slides, slide_map)

    if args.json:
        print(json.dumps(result, indent=2))
    else:
        if result["passed"]:
            print("PASS: PPTX validation passed")
            print(f"  Slides: {len(result['slides_info'])}")
            for slide in result["slides_info"]:
                title = (
                    slide["title"][:50] + "..."
                    if len(slide.get("title", "")) > 50
                    else slide.get("title", "(no title)")
                )
                print(f"  Slide {slide['index']}: {title} ({slide['shape_count']} shapes)")
        else:
            print("FAIL: PPTX validation failed")
            print(f"  Issues: {len(result['issues'])}")
            for issue in result["issues"]:
                print(f"  - {issue}")

    sys.exit(0 if result["passed"] else 2)


if __name__ == "__main__":
    main()
